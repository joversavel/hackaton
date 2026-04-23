from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Brand colours ──────────────────────────────────────────
NAVY       = RGBColor(0x14, 0x46, 0x8C)
NAVY_DARK  = RGBColor(0x0e, 0x32, 0x68)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG   = RGBColor(0xEE, 0xF2, 0xF8)
TEXT_DARK  = RGBColor(0x0F, 0x1F, 0x35)
ACCENT     = RGBColor(0x16, 0xA3, 0x4A)   # green
ORANGE     = RGBColor(0xD9, 0x77, 0x06)
SLATE      = RGBColor(0x33, 0x41, 0x55)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]  # blank

# ── helpers ────────────────────────────────────────────────

def add_rect(slide, left, top, width, height, fill_rgb, alpha=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    return shape

def add_text(slide, text, left, top, width, height,
             font_size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf    = txBox.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox

def add_bullet_box(slide, items, left, top, width, height,
                   font_size=16, color=TEXT_DARK, icon="•"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf    = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"{icon}  {item}"
        run.font.size  = Pt(font_size)
        run.font.color.rgb = color
    return txBox

def navy_header(slide, title, subtitle=None):
    """Full-width navy top bar."""
    bar_h = Inches(1.5)
    add_rect(slide, 0, 0, W, bar_h, NAVY_DARK)
    add_text(slide, title,
             Inches(0.5), Inches(0.18), Inches(12), Inches(0.8),
             font_size=32, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.5), Inches(0.95), Inches(12), Inches(0.45),
                 font_size=15, color=RGBColor(0xC5, 0xD5, 0xE9))

def accent_bar(slide):
    """Thin accent line below header."""
    add_rect(slide, 0, Inches(1.5), W, Inches(0.04), ACCENT)

# ══════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)

# Full background
add_rect(sl, 0, 0, W, H, NAVY_DARK)

# Big diagonal accent block
accent = sl.shapes.add_shape(1,
    Inches(7.8), 0, Inches(5.53), H)
accent.fill.solid()
accent.fill.fore_color.rgb = NAVY
accent.line.fill.background()

# White accent strip
add_rect(sl, Inches(7.6), 0, Inches(0.07), H, ACCENT)

# Title
add_text(sl, "Unified AI Workspace",
         Inches(0.6), Inches(1.8), Inches(7), Inches(1.1),
         font_size=42, bold=True, color=WHITE)
add_text(sl, "Assistant",
         Inches(0.6), Inches(2.8), Inches(7), Inches(0.9),
         font_size=42, bold=True, color=RGBColor(0xC5, 0xD5, 0xE9))

add_text(sl, "Microsoft Teams × Jira × Confluence × Dynamics 365",
         Inches(0.6), Inches(3.7), Inches(7), Inches(0.5),
         font_size=16, color=RGBColor(0x94, 0xA3, 0xB8))

add_text(sl, "AI Hackathon 2026  ·  Pollet Group IT",
         Inches(0.6), Inches(4.3), Inches(7), Inches(0.4),
         font_size=14, color=RGBColor(0x64, 0x74, 0x8B), italic=True)

# Right-side icons / labels
right_items = ["🤖  Claude AI (Anthropic)", "📋  Jira", "📚  Confluence", "👥  Azure AD", "💼  Dynamics 365"]
for i, item in enumerate(right_items):
    add_text(sl, item,
             Inches(8.5), Inches(1.8 + i * 0.9), Inches(4.5), Inches(0.6),
             font_size=17, bold=(i == 0), color=WHITE)

# Bottom bar
add_rect(sl, 0, Inches(7.1), W, Inches(0.4), NAVY)
add_text(sl, "Sterima  ·  2026",
         Inches(0.5), Inches(7.1), Inches(5), Inches(0.38),
         font_size=11, color=RGBColor(0x94, 0xA3, 0xB8))

# ══════════════════════════════════════════════════════════════
# SLIDE 2 — Het Probleem
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, LIGHT_BG)
navy_header(sl, "Het Probleem", "Waarom hebben we een AI-assistent nodig?")
accent_bar(sl)

problems = [
    ("⏱️  15–30 min per ticket",   "Support zoekt handmatig in Confluence naar bestaande oplossingen"),
    ("📋  Incomplete tickets",       "Developers missen context en richting bij Jira tickets"),
    ("🗂️  Versnipperde klantinfo",  "Dynamics, e-mail en tickets geven nooit een volledig beeld"),
    ("🧠  Kennisverliess",           "Opgeloste problemen worden niet gedocumenteerd voor hergebruik"),
]

for i, (title, desc) in enumerate(problems):
    top  = Inches(1.75 + i * 1.3)
    card = add_rect(sl, Inches(0.4), top, Inches(12.5), Inches(1.1),
                    RGBColor(0xFF, 0xFF, 0xFF))
    add_rect(sl, Inches(0.4), top, Inches(0.08), Inches(1.1), NAVY)
    add_text(sl, title,
             Inches(0.65), top + Inches(0.1), Inches(4), Inches(0.45),
             font_size=16, bold=True, color=TEXT_DARK)
    add_text(sl, desc,
             Inches(0.65), top + Inches(0.52), Inches(11.5), Inches(0.45),
             font_size=14, color=SLATE)

# ══════════════════════════════════════════════════════════════
# SLIDE 3 — De Oplossing
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, LIGHT_BG)
navy_header(sl, "De Oplossing", "Eén AI-interface voor alle systemen")
accent_bar(sl)

add_text(sl,
    "De Unified AI Workspace Assistant verbindt Jira, Confluence, Dynamics 365 en Azure AD "
    "in één Teams-interface — aangedreven door Claude (Anthropic).",
    Inches(0.5), Inches(1.65), Inches(12.3), Inches(0.8),
    font_size=15, color=SLATE)

cols = [
    ("🎯  Smart Assign",    ["Analyseert historiek", "Toewijzen met confidence %", "Directe Teams-notificatie"]),
    ("🔍  Oplossingszoeker", ["Vergelijkt met historiek", "Concrete aanpak", "Succeswaarschijnlijkheid"]),
    ("⚡  Auto-Resolve",    ["Analyseer + toewijzen", "Oplossing + comment", "Alles in één actie"]),
    ("📝  Kenniscapture",   ["Oplossing → artikel", "Auto-publish Confluence", "Kennis herbruikbaar"]),
]

col_w = Inches(3.1)
for i, (title, bullets) in enumerate(cols):
    left = Inches(0.3 + i * 3.25)
    top  = Inches(2.6)
    add_rect(sl, left, top, col_w, Inches(3.8), WHITE)
    add_rect(sl, left, top, col_w, Inches(0.5), NAVY)
    add_text(sl, title,
             left + Inches(0.1), top + Inches(0.07), col_w - Inches(0.15), Inches(0.4),
             font_size=13, bold=True, color=WHITE)
    add_bullet_box(sl, bullets,
                   left + Inches(0.15), top + Inches(0.65), col_w - Inches(0.25), Inches(2.9),
                   font_size=13, color=SLATE, icon="→")

# ══════════════════════════════════════════════════════════════
# SLIDE 4 — 5 Demo Scenario's
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, LIGHT_BG)
navy_header(sl, "5 Live Demo Scenario's", "De volledige demo flow")
accent_bar(sl)

scenarios = [
    ("1", "Jira Overview",     "Bot toont open tickets en vlagt onbeheerde tickets",          NAVY),
    ("2", "Ticket + Smart Assign", "Nieuw ticket → historiek → juiste persoon met confidence %", NAVY),
    ("3", "Oplossingszoeker",  "Claude vergelijkt met historiek → concrete aanpak + % succes",  NAVY),
    ("4", "Auto-Resolve",      "Analyseer, toewijzen, oplossing + statusupdate in één stap",   ACCENT),
    ("5", "Kenniscapture",     "Oplossing → kennisartikel auto-gepubliceerd in Confluence",    ORANGE),
]

for i, (num, title, desc, col) in enumerate(scenarios):
    top = Inches(1.7 + i * 1.02)
    add_rect(sl, Inches(0.35), top, Inches(0.6), Inches(0.82), col)
    add_text(sl, num,
             Inches(0.35), top + Inches(0.14), Inches(0.6), Inches(0.5),
             font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(sl, Inches(0.98), top, Inches(11.9), Inches(0.82), WHITE)
    add_text(sl, title,
             Inches(1.1), top + Inches(0.07), Inches(3.8), Inches(0.36),
             font_size=15, bold=True, color=TEXT_DARK)
    add_text(sl, desc,
             Inches(1.1), top + Inches(0.44), Inches(11.5), Inches(0.34),
             font_size=13, color=SLATE)

# ══════════════════════════════════════════════════════════════
# SLIDE 5 — Smart Assign Logic
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, LIGHT_BG)
navy_header(sl, "Smart Assign Logic", "Claude analyseert elk nieuw ticket")
accent_bar(sl)

# Left: inputs
add_rect(sl, Inches(0.35), Inches(1.65), Inches(5.5), Inches(5.6), WHITE)
add_rect(sl, Inches(0.35), Inches(1.65), Inches(5.5), Inches(0.42), NAVY)
add_text(sl, "📥  Input factoren",
         Inches(0.5), Inches(1.7), Inches(5), Inches(0.35),
         font_size=14, bold=True, color=WHITE)
inputs = [
    "Ticketinhoud, component & fouttype",
    "Historiek: wie loste dit eerder op?",
    "Gemiddelde oplostijd per persoon",
    "Huidige werkdruk & beschikbaarheid",
]
add_bullet_box(sl, inputs,
               Inches(0.5), Inches(2.2), Inches(5.1), Inches(4.5),
               font_size=14, color=SLATE, icon="▸")

# Arrow
add_text(sl, "→",
         Inches(6.05), Inches(3.6), Inches(0.8), Inches(0.6),
         font_size=36, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# Right: outputs
add_rect(sl, Inches(7.0), Inches(1.65), Inches(5.95), Inches(5.6), WHITE)
add_rect(sl, Inches(7.0), Inches(1.65), Inches(5.95), Inches(0.42), ACCENT)
add_text(sl, "📤  Output per ticket",
         Inches(7.15), Inches(1.7), Inches(5.5), Inches(0.35),
         font_size=14, bold=True, color=WHITE)
outputs = [
    "Beste kandidaat met confidence %",
    "Alternatief indien 1e keuze niet beschikbaar",
    "Kant-en-klare oplossingsvoorstel als comment",
    "Directe Teams-notificatie aan de toegewezene",
    "Uitleg van de redenering (niet alleen het getal)",
]
add_bullet_box(sl, outputs,
               Inches(7.15), Inches(2.2), Inches(5.6), Inches(4.5),
               font_size=14, color=SLATE, icon="✓")

# ══════════════════════════════════════════════════════════════
# SLIDE 6 — Technische Stack
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, LIGHT_BG)
navy_header(sl, "Technische Stack", "Bewezen enterprise technologieën")
accent_bar(sl)

stack = [
    ("🤖  AI Engine",        "Claude API — Anthropic\nclaude-sonnet-4-6 / claude-opus-4-7 + prompt caching"),
    ("📋  Tickets & Docs",   "Atlassian MCP — Jira + Confluence (native read/write)"),
    ("💬  Chat Interface",   "Microsoft Teams Bot Framework"),
    ("🔐  Authenticatie",    "Azure AD — bestaande infrastructuur"),
    ("💼  CRM",              "Dynamics 365 connector (of mock data tijdens demo)"),
    ("🌐  Frontend",         "Flask web app — light/dark mode, Pollet brand kit"),
]

col2_start = 3
for i, (comp, tech) in enumerate(stack):
    col   = i // col2_start
    row   = i %  col2_start
    left  = Inches(0.35 + col * 6.55)
    top   = Inches(1.75 + row * 1.75)
    add_rect(sl, left, top, Inches(6.3), Inches(1.55), WHITE)
    add_rect(sl, left, top, Inches(0.06), Inches(1.55), NAVY)
    add_text(sl, comp,
             left + Inches(0.2), top + Inches(0.12), Inches(5.8), Inches(0.4),
             font_size=14, bold=True, color=TEXT_DARK)
    add_text(sl, tech,
             left + Inches(0.2), top + Inches(0.55), Inches(5.9), Inches(0.85),
             font_size=12, color=SLATE)

# ══════════════════════════════════════════════════════════════
# SLIDE 7 — Business Value
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, LIGHT_BG)
navy_header(sl, "Business Value", "Voor elke afdeling een direct voordeel")
accent_bar(sl)

depts = [
    ("Support",        "Snellere responstijden, minder handmatig zoeken",          NAVY),
    ("Development",    "Volledige context + oplossingsvoorstel bij elk ticket",     NAVY),
    ("CRM / Sales",    "360° klantoverzicht in seconden",                           ACCENT),
    ("IT / Azure",     "Native Azure AD + MCP integratie, geen extra infra",        ACCENT),
    ("Management",     "Minder silo's, kortere doorlooptijden",                    ORANGE),
    ("Iedereen",       "Kennis bewaard en herbruikbaar — geen kennisverliess meer", ORANGE),
]

col_count = 3
for i, (dept, value, col_color) in enumerate(depts):
    col  = i % col_count
    row  = i // col_count
    left = Inches(0.3  + col * 4.35)
    top  = Inches(1.75 + row * 2.55)
    add_rect(sl, left, top, Inches(4.1), Inches(2.3), WHITE)
    add_rect(sl, left, top, Inches(4.1), Inches(0.45), col_color)
    add_text(sl, dept,
             left + Inches(0.15), top + Inches(0.07), Inches(3.8), Inches(0.34),
             font_size=14, bold=True, color=WHITE)
    add_text(sl, value,
             left + Inches(0.15), top + Inches(0.6), Inches(3.8), Inches(1.5),
             font_size=13, color=SLATE)

# ══════════════════════════════════════════════════════════════
# SLIDE 8 — Closing / Next Steps
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, NAVY_DARK)
add_rect(sl, 0, 0, Inches(0.12), H, ACCENT)

add_text(sl, "Klaar voor de demo.",
         Inches(0.5), Inches(1.5), Inches(9), Inches(1.1),
         font_size=44, bold=True, color=WHITE)

add_text(sl, "Unified AI Workspace Assistant — AI Hackathon 2026",
         Inches(0.5), Inches(2.65), Inches(9), Inches(0.5),
         font_size=16, color=RGBColor(0x94, 0xA3, 0xB8), italic=True)

next_steps = [
    "5 live demo scenario's doorlopen",
    "Smart Assign confidence % uitleg",
    "Auto-Resolve in real-time",
    "Kennisartikel gepubliceerd in Confluence",
]
add_bullet_box(sl, next_steps,
               Inches(0.5), Inches(3.4), Inches(8), Inches(2.5),
               font_size=18, color=WHITE, icon="→")

add_text(sl, "Pollet Group IT  ·  Sterima  ·  2026",
         Inches(0.5), Inches(6.9), Inches(8), Inches(0.4),
         font_size=12, color=RGBColor(0x64, 0x74, 0x8B))

# ── Save ───────────────────────────────────────────────────
out = "Pollet_AI_Workspace_Assistant.pptx"
prs.save(out)
print(f"Saved: {out}")
